from pptx import Presentation
from pptx.util import Inches, Pt
from os import listdir
from os.path import isfile
import pandas as pd

if isfile("./graphs/previous_clusters.csv"):
    previousAnnotation = True
    previousAnnData = pd.read_csv("./graphs/previous_clusters.csv", index_col = 0 )
else:
    previousAnnotation = False

# count the number of clusters by counting the number of images begin with "clusters"
clusters = len([tsne_img for tsne_img in listdir("./graphs") if tsne_img[0:7] == "cluster"])

# initiate presentation
prs = Presentation()

prs = Presentation()
black_slide_layout = prs.slide_layouts[6]

# add tsne and umap plots by cluster on first slide
img_path = "./graphs/dr.png"
slide = prs.slides.add_slide(black_slide_layout)
left = Inches(.1)
top = Inches(.5)
height = Inches(5.7)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

# add tsne and umap plots by sample on second slide
img_path = "./graphs/dr_sample.png"
slide = prs.slides.add_slide(black_slide_layout)
left = Inches(.1)
top = Inches(.5)
height = Inches(5.7)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

# for each cluter counted import dr plot and tally plots and insert them on a slide
for cluster in range(clusters):
    # insert dr plot
    img_path = "./graphs/cluster_dr_{cluster}.png".format(cluster = cluster)
    slide = prs.slides.add_slide(black_slide_layout)
    left = Inches(7)
    top = Inches(.7)
    height = Inches(6.7)
    pic = slide.shapes.add_picture(img_path, left, top, height=height)
    # insert tally plot
    img_path = "./graphs/tally_{cluster}.png".format(cluster = cluster)
    left = Inches(2.2)
    top = Inches(0)
    height = Inches(.77)
    pic = slide.shapes.add_picture(img_path, left, top, height=height)
    # insert the text
    left = top = Inches(0)
    width = Inches(6)
    height = Inches(3)
    txtBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txtBox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Cluster {cluster}: ...".format(cluster = int(cluster))
    p.font.bold = True
    p.font.size = Pt(24)
    p = tf.add_paragraph()
    p.text = "Defining markers:"
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "..."; p.level = 1
    p.font.size = Pt(10)
    p = tf.add_paragraph(); p.text = "Indentity: ..."; p.font.bold = True;
    p.font.size = Pt(12)
    p = tf.add_paragraph(); p.text = "Justification: ..."; p.font.bold = True;
    p.font.size = Pt(12)

    if previousAnnotation:
        left = Inches(4); top = Inches(1); width = Inches(2); height = Inches(3);
        txtBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txtBox.text_frame; tf.clear(); p = tf.paragraphs[0];
        p.text = "\n".join(previousAnnData.loc[cluster].values[0].split("; "))
        p.font.size = Pt(12)
        
    
prs.save('annotation_template.pptx')
